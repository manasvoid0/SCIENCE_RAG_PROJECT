"""
text_extractor.py
=================
Job: take an uploaded file (PDF / Word / PowerPoint) and pull out ONLY the words.
Images, fonts, colours and layout are all thrown away -- we just want raw text
so the detectors can analyse it.

Each file type needs a different library because each format stores text differently.
"""

import os  # used to look at the file extension (.pdf, .docx, .pptx)

# NOTE: the PDF/Word/PPT libraries are imported INSIDE each function below
# (not at the top) so that plain .txt files work even before you run
# `pip install -r requirements.txt`.


def _extract_pdf(filepath: str) -> str:
    """Open a PDF and join the text of every page into one big string."""
    import fitz  # PyMuPDF (install name: pymupdf)
    doc = fitz.open(filepath)  # open the PDF
    pages = [page.get_text() for page in doc]  # grab text from each page
    doc.close()  # always close the file when done
    return " ".join(pages)  # glue all pages together with spaces


def _extract_docx(filepath: str) -> str:
    """Open a Word file and join the text of every paragraph."""
    import docx  # python-docx
    document = docx.Document(filepath)  # open the .docx
    paragraphs = [p.text for p in document.paragraphs]  # one string per paragraph
    return " ".join(paragraphs)


def _extract_pptx(filepath: str) -> str:
    """Open a PowerPoint and join the text from every shape on every slide."""
    from pptx import Presentation  # python-pptx
    prs = Presentation(filepath)  # open the .pptx
    chunks = []  # collect text pieces here
    for slide in prs.slides:  # walk through each slide
        for shape in slide.shapes:  # walk through each box/shape on the slide
            if shape.has_text_frame:  # only shapes that actually contain text
                chunks.append(shape.text_frame.text)
    return " ".join(chunks)


def extract_text(filepath: str) -> str:
    """
    The ONE function the rest of the app calls.
    It looks at the file's extension and routes to the correct reader above.
    """
    lower = filepath.lower()  # lowercase so ".PDF" and ".pdf" both work

    if lower.endswith(".txt"):
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()  # plain text needs no special library
    elif lower.endswith(".pdf"):
        text = _extract_pdf(filepath)
    elif lower.endswith(".docx"):
        text = _extract_docx(filepath)
    elif lower.endswith(".pptx"):
        text = _extract_pptx(filepath)
    else:
        # If someone uploads an unsupported type, fail loudly with a clear message
        raise ValueError(f"Unsupported file type: {os.path.basename(filepath)}")

    # Normalise whitespace: collapse newlines/tabs/multiple spaces into single spaces
    return " ".join(text.split())
